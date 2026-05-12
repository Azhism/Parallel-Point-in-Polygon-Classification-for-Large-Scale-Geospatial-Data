"""
Updates PDC_presentation.pptx:
  1. Inserts a Team Contributions slide after slide 5 (Project Roadmap)
  2. Fixes M2 Results slide: updates "2.14x" to "2.16x" and improves label
  3. Adds a brief note to M3 results confirming multi-process worker path
"""
import sys
import copy
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colors (matching existing slide theme) ────────────────────────────────────
BG        = RGBColor(0x0A, 0x16, 0x28)   # #0A1628 dark navy
CYAN      = RGBColor(0x22, 0xD3, 0xEE)   # #22D3EE accent
WHITE     = RGBColor(0xF1, 0xF5, 0xF9)   # #F1F5F9 main text
SLATE     = RGBColor(0x94, 0xA3, 0xB8)   # #94A3B8 subtext
CARD_BG   = RGBColor(0x10, 0x24, 0x3E)   # slightly lighter navy for cards
CARD_BD   = RGBColor(0x22, 0xD3, 0xEE)   # cyan border

# ── Slide dimensions (20" × 11.25") ───────────────────────────────────────────
W = Inches(20)
H = Inches(11.25)


def inches(val):
    return Inches(val)


def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(inches(left), inches(top),
                                     inches(width), inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None,
             border_width_pt=1.0):
    from pptx.util import Pt as PtU
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        inches(left), inches(top), inches(width), inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width_pt)
    else:
        shape.line.fill.background()
    return shape


def set_slide_background(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def move_slide(prs, from_index, to_index):
    """Move slide at from_index to to_index (0-based)."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    elem = slides_list[from_index]
    xml_slides.remove(elem)
    if to_index >= len(list(xml_slides)):
        xml_slides.append(elem)
    else:
        ref = list(xml_slides)[to_index]
        xml_slides.insert(list(xml_slides).index(ref), elem)


def build_team_contributions_slide(prs):
    """Create and return a Team Contributions slide."""
    blank_layout = prs.slide_layouts[6]   # blank layout
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BG)

    # ── Header ────────────────────────────────────────────────────────────────
    add_textbox(slide, 1.67, 1.10, 6.0, 0.35,
                "TEAM CONTRIBUTIONS", 22, bold=True, color=CYAN)
    add_textbox(slide, 1.67, 1.55, 17.0, 0.75,
                "Who built what, milestone by milestone", 28, color=WHITE)

    # thin separator line (narrow rect)
    add_rect(slide, 1.67, 2.50, 16.66, 0.03, CYAN)

    # ── Four member cards ─────────────────────────────────────────────────────
    members = [
        {
            "name":  "Azhab Safwan Babar",
            "erp":   "ERP 29227",
            "m1":    "M1: Quadtree spatial index + bounding-box filtering",
            "m2":    "M2: Spatial partitioning (grid / hierarchical tiles)",
            "m3":    "M3: Distributed execution model (multi-process IPC)",
        },
        {
            "name":  "Ayaan Malik",
            "erp":   "ERP 29238",
            "m1":    "M1: Ray-casting PiP algorithm (edge, vertex, holes)",
            "m2":    "M2: Multithreading setup — OpenMP static & dynamic",
            "m3":    "M3: Batch processing pipeline (streaming 250K chunks)",
        },
        {
            "name":  "Muhammad Anas Tabba",
            "erp":   "ERP 28998",
            "m1":    "M1: Dataset generation — uniform & clustered GPS pts",
            "m2":    "M2: Dynamic task queues & work-stealing deques",
            "m3":    "M3: Polygon replication vs spatial sharding trade-off",
        },
        {
            "name":  "Tayyab Bin Hisham",
            "erp":   "ERP 29087",
            "m1":    "M1: Correctness testing & benchmark validation",
            "m2":    "M2: Cache locality — Morton/Z-order tiled sorting",
            "m3":    "M3: Result aggregation (count / checksum per worker)",
        },
    ]

    # card geometry — two rows, two columns
    col_x   = [1.10, 10.55]   # left edges
    row_y   = [2.70, 6.85]    # top edges
    c_w     = 8.60             # card width
    c_h     = 3.80             # card height

    positions = [
        (col_x[0], row_y[0]),
        (col_x[1], row_y[0]),
        (col_x[0], row_y[1]),
        (col_x[1], row_y[1]),
    ]

    for (lx, ty), m in zip(positions, members):
        # card background rectangle
        add_rect(slide, lx, ty, c_w, c_h, CARD_BG, CARD_BD, 0.8)

        # member name (cyan, large)
        add_textbox(slide, lx + 0.22, ty + 0.20, c_w - 0.44, 0.50,
                    m["name"], 24, bold=True, color=CYAN)
        # ERP number
        add_textbox(slide, lx + 0.22, ty + 0.72, c_w - 0.44, 0.30,
                    m["erp"], 16, color=SLATE)
        # milestone contributions
        add_textbox(slide, lx + 0.22, ty + 1.12, c_w - 0.44, 0.42,
                    m["m1"], 15, color=WHITE)
        add_textbox(slide, lx + 0.22, ty + 1.58, c_w - 0.44, 0.42,
                    m["m2"], 15, color=WHITE)
        add_textbox(slide, lx + 0.22, ty + 2.04, c_w - 0.44, 0.42,
                    m["m3"], 15, color=WHITE)

    # bottom note
    add_textbox(slide, 1.10, 10.75, 17.80, 0.35,
                "All members collaborated on performance benchmarking and scaling analysis.",
                14, color=SLATE, align=PP_ALIGN.CENTER)

    return slide


# ── Main update logic ──────────────────────────────────────────────────────────
INPUT  = "d:/Classess/PDC/Project/PDC_presentation.pptx"
OUTPUT = "d:/Classess/PDC/Project/PDC_presentation.pptx"

prs = Presentation(INPUT)

print(f"Loaded presentation: {len(prs.slides)} slides")

# ────────────────────────────────────────────────────────────────────────────
# 1. Fix M2 Results slide (currently slide 12, index 11)
#    Change "2.14x" -> "2.16x" (canonical Hybrid speedup for 1M uniform)
#    and update the label for clarity.
# ────────────────────────────────────────────────────────────────────────────
m2_slide = prs.slides[11]
for shape in m2_slide.shapes:
    if not hasattr(shape, "text_frame"):
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip() == "2.14×":
                run.text = "2.16×"
                print("  Fixed: '2.14x' -> '2.16x' on M2 results slide")
            elif run.text.strip() == "SPEEDUP @ 4 THREADS · 1M PTS":
                run.text = "HYBRID SPEEDUP · 1M UNIFORM"
                print("  Fixed: M2 label updated")

# ────────────────────────────────────────────────────────────────────────────
# 2. Add Team Contributions slide (will be appended, then moved to index 5)
# ────────────────────────────────────────────────────────────────────────────
print("Building Team Contributions slide...")
build_team_contributions_slide(prs)
new_slide_index = len(prs.slides) - 1   # currently last slide (index 18)
print(f"  Slide added at index {new_slide_index} (total: {len(prs.slides)})")

move_slide(prs, new_slide_index, 5)   # insert after slide 5 (roadmap, index 4)
print(f"  Slide moved to index 5")

# ────────────────────────────────────────────────────────────────────────────
# 3. Save
# ────────────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\nSaved: {OUTPUT}")
print(f"Final slide count: {len(prs.slides)}")
